#Import all packages
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from selenium import webdriver
from mutagen.mp3 import MP3
import moviepy.editor as mp
from PIL import Image, ImageEnhance
import boto3
import glob
import os

#To do list:

#DO BEFORE FIRST PRODUCTION
#  in compile():
#    a: add music 
#  speed up audio
#  moviepy takes too long, do this shit with ffmpeg
#  add nuerel voice to aws

#DO BEFORE REGULAR PRODUCTION
#  add more background vid options
#  Make image conversion automatic?
#  make it so multiple videos can be queued at once
#research youtube algorithm/upload schedules

#DO BEFORE PUBLISH
#  Clean up comments
#  add input systems
#  add command line support
#  clean up print statements
#  make custom resize values depending on vid
#  make text on slides smaller and take up more space per line
#  add def split_text() (Saved in YouTubeFucker)?

class Comments:
    audio_list = []
    cwd=os.getcwd()
    #PATH_TO_AUDIO_READY_FILE="/Users/egray/Desktop/YoutubeFucker/RedditCreator/List_of_comments/audio_ready_test.txt" # Create program that translates them
    #COMMENT_FILE=cwd+"/comment_list.txt"
    AUDIO_START="aws_audio_"
    DOWNLOAD_PATH="/Users/egray/Downloads/"
    CONVERTIO_URL="https://convertio.co/pptx-png/"
    GECKO_PATH=cwd+"/geckodriver"
    #comments_file=open(cwd+"/comment_list.txt", "r")
    dir=cwd+"/"
    comment_count=0
    length_of_clip=0
    slide_count=0
    pic_num=0
    count_list=[]
    durations=[]
    clip_amount=[]
    video_list=[]
    cleanup_list=[]
    image_path=""
    image_start="image"
    clipLenDictionary={}
    def __init__(self, output_name, accent, title, background_choice):
        self.output_name = output_name #  Name of the compiled audio clip
        #self.dir = dir #  Path to current directory
        self.accent = accent #  AWS accent
        self.title = title #  If i need to explain this to myself, i should give up coding. I need sleep.
        self.background_choice = background_choice #The background file i picked

    
    def text2slide(self):
        #Open DOC
        print("Converting text to slides...")
        comments_file=open(self.dir+"comment_list.txt", "r")

        lines=comments_file.readlines()

        
        #Var decelartion for layout
        prs=Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        left = top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)

        t_left = Inches(0.5)
        t_top = Inches(0)
        t_width = Inches(9)
        t_height = Inches(1)

        i=0 #counts how many paragraphs are in one slide
        count=0 #counts total number of slides
        trans_count=0 #number of transitions/comments that start with u/
        paragraph_length=0 
        
        paragraph_text=""
        new_title=True
        firstComment=True
        for line in lines:
            if(line.startswith("u/")):
                user=line
                paragraph_length=0
                paragraph_text=""
                #adds string to dictionary for transition times
                if firstComment==False:
                    clip_str="transition"+str(trans_count)
                    trans_count+=1
                    self.clipLenDictionary.update({clip_str : 2.0})
                else:
                    firstComment=False

                
                i=0
                
            else:
                i=i+1
                clip_str="clip"+str(count)
                count=count+1
                self.clipLenDictionary.update({clip_str : "duration"})
                slide = prs.slides.add_slide(blank_slide_layout)
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(34, 34, 34)

                txBox = slide.shapes.add_textbox(left, top, width, height)
                paragraph_frame = txBox.text_frame
                userBox = slide.shapes.add_textbox(t_left, t_top, t_width, t_height)
                user_frame = userBox.text_frame
                user_text = user_frame.add_paragraph()

                user_text.font.size=Pt(17)
                user_text.font.name='Verdana'
                user_text.font.color.rgb = RGBColor(20, 158, 240)
                user_text.text=user
                
                paragraph_frame.word_wrap=True
                paragraph=paragraph_frame.add_paragraph()
                paragraph.font.color.rgb = RGBColor(239, 239, 237)
                paragraph.font.size=Pt(16)
                paragraph.font.name="Verdana"
                line_length=len(line)
                paragraph_length=paragraph_length+line_length
                if paragraph_length<1200:
                    if i==1:
                        paragraph_text=paragraph_text+line
                        paragraph_text=paragraph_text+"\n \n"
                        paragraph.text=paragraph_text
                    if i==2:
                        paragraph_text=paragraph_text+line
                        paragraph_text=paragraph_text+"\n \n"
                        paragraph.text=paragraph_text
                    if i==3:
                        paragraph_text=paragraph_text+line
                        paragraph_text=paragraph_text+"\n \n"
                        paragraph.text=paragraph_text
                    if i==4:
                        paragraph_text=paragraph_text+line
                        paragraph_text=paragraph_text+"\n \n"
                        paragraph.text=paragraph_text
                        i=0
                        paragraph_text=""
                else:
                    i=0
                    paragraph_text=""
                    paragraph_text=paragraph_text+line
                    paragraph_text=paragraph_text+"\n \n"
                    paragraph.text=paragraph_text
                    paragraph_length=len(paragraph_text)

        powerpoint_name=self.dir+self.title+".pptx"
        self.cleanup_list.append(self.title+".pptx")

        prs.save(powerpoint_name) # saving file
        comments_file.close()
        print("Powerpoint created!")

    
    #rewrite text2aduio
    def create_clips(self):
        comments_file=open(self.cwd+"/comment_list.txt", "r")
        lines=comments_file.readlines()
        i=0
        clip_count=0
        duration=0
        first=True
        for line in lines:
            if(line.startswith("u/")):
                self.comment_count=self.comment_count+1
                self.count_list.append(self.slide_count)
                self.slide_count=0
                
                if(first==False):
                    self.audio_list.append("/Users/egray/Desktop/YoutubeFucker/RedditCreator/List_of_comments/transition.mp3")
                    
                else:
                    first=False

            else:
                file_name=self.dir+self.AUDIO_START+str(i)+".mp3"
                #audio_str="<speak>"+line+"<break strength=\"strong\"/></speak>"
                self.create_clip(line, file_name, i)
                self.audio_list.append(file_name)
                self.audio_list.append("/Users/egray/Desktop/RedditMonkey/pause.mp3")
                self.cleanup_list.append(self.AUDIO_START+str(i)+".mp3")
                self.slide_count=self.slide_count+1
                audio=MP3(file_name)
                audio_length=audio.info.length
                clip_str="clip"+str(clip_count)
                #print(clip_str)
                self.clipLenDictionary[clip_str] = audio_length+1
                
                clip_count+=1
                i+=1
        clips=0
        for name, length in self.clipLenDictionary.items():
            if name.startswith("t"):
                self.durations.append(duration)
                duration=0
                self.clip_amount.append(clips)
                clips=0
            else:
                duration=duration+length
                clips+=1
        self.durations.append(duration)
        self.clip_amount.append(clips)
        self.audio_list.append(file_name)
        #print(self.clipLenDictionary)
        #print("\n\n")
        #print(self.durations)
        #print(self.clip_amount)
        #print("AUDIO LIST")
        comments_file.close()
        #return self.audio_list



    def create_clip(self, text, filename, i):
        client = boto3.client("polly")
        #if self.accent=="uk" or self.accent=="UK":
            #response=client.synthesize_speech(Text=text, LanguageCode="en-GB", VoiceId='Brian', OutputFormat="mp3") #British

        #elif self.accent=="us" or self.accent=="US":
        response=client.synthesize_speech(Text=text, VoiceId='Matthew', OutputFormat="mp3") #American

        body = response["AudioStream"].read()

        with open(filename, 'wb') as file:
            file.write(body)
            #print("Writting "+self.AUDIO_START+str(i))
            file.close()
    def convert_pptx(self):
        driver = webdriver.Firefox(executable_path = self.GECKO_PATH)
        driver.get(self.CONVERTIO_URL)

        print("Add the powerpoint file to convertio, then click download")
        input("Press enter when ready...  ")
        driver.quit()
        print("Quiting...")
        print("Done!")

        guess=self.DOWNLOAD_PATH+self.title+".png/"
        print("Printing guess at new file ")
        print(guess)
        user_input=input("Is this the correct name for the newly downloaded file? [Y/N]:  ")
        if user_input=="y" or user_input=="Y":
            print("FUCK YEAH")
            self.image_path=guess
        elif user_input=="n" or user_input=="N":
            print("SHIT")
            new_path=input("What is the correct file name?  ")
            self.image_path=self.DOWNLOAD_PATH+new_path+"/"

        glob_str=self.image_path+"*.png"
        image_list=glob.glob(glob_str)
        self.pic_num=len(image_list)
        #print(image_list)
        for item in image_list:
            if item[-6:]=="-0.png":
                self.image_start=item[0:-6]
                break
        for item in image_list: # TO DO: CHANGE OPACITY OF EVERY IMAGE
            image=Image.open(item)
            image=image.convert("RGBA")
            image.putalpha(190)
            image.save(item, "PNG")

    def concat_audios(self):
        st = ["ffmpeg -loglevel quiet"]#add 
        print("Compiling audio...\n")
        #print(len(self.audio_list))
        for i in range(len(self.audio_list)):
            st.append("-i")
            st.append(self.audio_list[i])

        st.append("-filter_complex")
        st = ' '.join(st) + " "

        for i in range(len(self.audio_list)):
            st = st + "[" + str(i) + ":0]"

        st = st + "concat=n=" + str(len(self.audio_list)) + ":v=0:a=1[out] -map [out] " + self.dir + self.output_name
        os.system(st)

        print("Audio compiled!...\n")

    def videoMaker(self):
        count=0
        total=0
        vid_number=0
        i=1
        it=0
        temp_sounds=[]
        first=True
        first_audio=self.dir+"aws_audio_"+str(total)+".mp3"
        for duration in self.durations:

            #good, strong, communist working code
            clip_count=self.clip_amount[vid_number]
            vid_number+=1
            #print("Clip count: "+str(clip_count))
            first_audio=self.dir+"aws_audio_"+str(total)+".mp3"
            first_slide_name=self.image_start+"-"+str(total)+".png"
            first_clip_name="clip"+str(total)
            first_clip_len=self.clipLenDictionary.get(first_clip_name)
            #print("first slide name, first audio name, first clip len: ")
            print(first_slide_name, first_audio, first_clip_len)
            #print("\nprinting the rest...")
            total=total+clip_count
            time=0.000 #used to decide when clip should appear

            if duration > 0 and duration <= 30:
                background_clip = mp.VideoFileClip(self.background_choice+"_30s.mp4", audio=False)
            
            #   background_clip = mp.VideoFileClip("/Users/egray/Desktop/YoutubeFucker/RedditCreator/Reddit_stock/america_1minute.mp4")
            elif duration > 30 and duration <= 60:
                background_clip = mp.VideoFileClip(self.background_choice+"_60s.mp4", audio=False)
            #   background_clip = mp.VideoFileClip("/Users/egray/Desktop/redditTool/america_3min.mp4")
            elif duration > 60 and duration <= 300:
                background_clip = mp.VideoFileClip(self.background_choice+"_5m.mp4", audio=False)
            elif duration > 300 and duration <=600:
                background_clip = mp.VideoFileClip(self.background_choice+"_10m.mp4", audio=False)
            elif duration > 600 and duration <=900:
                background_clip = mp.VideoFileClip(self.background_choice+"_15m.mp4", audio=False)
            else:
                print("CLIP TOO LONG")

            #Terrible, unedited capatilist code
            first_slide=(mp.ImageClip(first_slide_name))
            first_slide=first_slide.resize(0.7)
            first_slide=first_slide.set_duration(first_clip_len) #sets how long pic should be up

            video = mp.CompositeVideoClip([background_clip, first_slide.set_position("center").set_start(time)])
            time+=first_clip_len 
            if(first==False):
                i+=1
            while i<total:
                audio_clip=self.dir+self.AUDIO_START+str(i)+".mp3"
                clip_name=self.image_start+"-"+str(i)+".png"
                #print("Next clip:\n")
                #print(audio_clip, clip_name)
                #print("Duration:")
                search_str="clip"+str(i)
                clip_len=self.clipLenDictionary.get(search_str)
                #print(clip_len)

                slide = (mp.ImageClip(clip_name))
                slide=slide.resize(0.7)
                slide=slide.set_duration(clip_len) #sets how long pic should be up
                video = mp.CompositeVideoClip([video, slide.set_position("center").set_start(time)])
                time=time+clip_len
                i+=1
            first=False
            video=video.subclip(0, duration)
            #audioclip = mp.AudioFileClip(self.output_name)

            #new_audioclip = mp.CompositeAudioClip([audioclip])
            #video.audio = new_audioclip
            vid_title=self.title+str(it)+".mp4"
            video.write_videofile(vid_title, audio=False)
            print("FINISHED WRITING:"+vid_title)
            #self.videos.append(vid_title)
            self.cleanup_list.append(vid_title)
            self.video_list.append(self.dir+vid_title)
            self.video_list.append(self.dir+"transition2.mp4")
            
            it=it+1
    #def add_image(self, img, clip_duration, start_time, video):
    def compile(self):
        list_len=len(self.video_list)-1
        i=0
        
        video_list_file=open("video_list.txt", "r+")
        for vid in self.video_list:
            if(i!=list_len):
                str="file \'"+vid+"\'\n"
                video_list_file.write(str)
            i=i+1
        video_list_file.close()
        os.system('ffmpeg -loglevel quiet -f concat -safe 0 -i video_list.txt -c copy video_noaudio.mp4')
        
        os_str='ffmpeg -loglevel quiet -i video_noaudio.mp4 -i '+self.output_name+' -map 0:v -map 1:a -c:v copy -shortest ' +self.title+'_FINALE.mp4'
        os.system(os_str)
        self.cleanup_list.append("video_noaudio.mp4")

        #os.system('ffmpeg -i video_noaudio.mp4 -i audio_output.mp3 -map 0:v -map 1:a -c:v copy -shortest FINALEMAYBE.mp4')
    def clean_up(self):
        #self.cleanup_list.append('geckodriver.log')
        #input("Press enter")
        for item in self.cleanup_list:
            os.remove(item)
        #os.remove(self.title+".png", dir_fd=self.DOWNLOAD_PATH)
        print("---FINISHED---")


#USAGE:  Comments(<output_name.mp3> <accent> <video_title>, <background_vid>)
p1 = Comments("pizza1_output.mp3", "US", "Pizza1", "/Users/egray/Desktop/YoutubeFucker/RedditCreator/Reddit_stock/Chill/chill2")


p1.text2slide()
p1.create_clips()
p1.concat_audios()
p1.convert_pptx()
p1.videoMaker()
p1.compile()
p1.clean_up()